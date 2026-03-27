#!/usr/bin/env python3
"""
Doc Creator - Generate DOCX/XLSX/PPTX documents

Supports:
- DOCX: Word documents with images, formatting, tables
- XLSX: Excel spreadsheets with data and charts
- PPTX: PowerPoint presentations with slides
"""

import argparse
import sys
from pathlib import Path
from typing import Optional

try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
except ImportError:
    print("Error: python-docx not installed. Run: pip install python-docx>=1.0.0")
    sys.exit(1)

import platform


def _get_system_cjk_font() -> str:
    """Return a CJK font name available on the current OS."""
    os_name = platform.system()
    if os_name == "Windows":
        return "Microsoft YaHei"
    elif os_name == "Darwin":
        return "PingFang SC"
    else:
        return "Noto Sans CJK SC"


def _set_chinese_font(doc: Document, font_name: str | None = None) -> None:
    """Set East Asian font for all styles to support CJK characters."""
    font = font_name or _get_system_cjk_font()
    # Set Normal style
    style_normal = doc.styles["Normal"]
    style_normal.font.name = font
    style_normal.element.rPr.rFonts.set(qn("w:eastAsia"), font)
    # Set all Heading styles
    for level in range(1, 5):
        key = f"Heading {level}"
        if key in doc.styles:
            h_style = doc.styles[key]
            h_style.font.name = font
            h_style.element.rPr.rFonts.set(qn("w:eastAsia"), font)

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill
    from openpyxl.utils import get_column_letter
except ImportError:
    print("Error: openpyxl not installed. Run: pip install openpyxl>=3.1.0")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx>=0.6.0")
    sys.exit(1)


def _add_image(doc, image_path: str) -> bool:
    """Add an image to the document, centered."""
    img_path = Path(image_path)
    if not img_path.exists():
        print(f"Warning: Image file not found: {image_path}")
        return False
    try:
        doc.add_picture(str(img_path), width=Inches(6))
        last_paragraph = doc.paragraphs[-1]
        last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        return True
    except Exception as e:
        print(f"Warning: Could not insert image {image_path}: {e}")
        return False


def create_docx(
    title: str,
    content: str,
    image_paths: list[str] | None = None,
    output_path: str = "output.docx",
    font_name: str | None = None,
) -> str:
    """
    Create a Word document with title, content, and optional images.

    Images can be placed inline in content using [IMAGE:N] markers (0-indexed).
    Images without markers are appended at the end.

    Args:
        title: Document title
        content: Main content text (supports ## H1, ### H2, - bullets, [IMAGE:N])
        image_paths: List of image file paths
        output_path: Output file path
        font_name: Override CJK font (auto-detected by default)

    Returns:
        Path to created document
    """
    # Ensure output path has .docx extension
    if not output_path.endswith('.docx'):
        output_path = f"{output_path}.docx"

    images = image_paths or []
    used_images = set()

    doc = Document()
    _set_chinese_font(doc, font_name=font_name)

    # Add title
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add content with inline image support
    if content:
        paragraphs = content.split('\n\n')
        for para_text in paragraphs:
            text = para_text.strip()
            if not text:
                continue

            # Check for inline image marker [IMAGE:N]
            import re
            img_marker = re.match(r'^\[IMAGE:(\d+)\]$', text)
            if img_marker:
                idx = int(img_marker.group(1))
                if idx < len(images):
                    _add_image(doc, images[idx])
                    used_images.add(idx)
                continue

            # Check if it's a heading (starts with #)
            if text.startswith('###'):
                doc.add_heading(text.replace('###', '').strip(), 2)
            elif text.startswith('##'):
                doc.add_heading(text.replace('##', '').strip(), 1)
            elif text.startswith('- '):
                # Bullet list - collect all consecutive bullet items
                items = [text[2:]]
                # (single paragraph for now; consecutive bullets are separate paragraphs)
                doc.add_paragraph(items[0], style='List Bullet')
            else:
                doc.add_paragraph(text)

    # Append any remaining images that weren't placed inline
    for idx, img in enumerate(images):
        if idx not in used_images:
            # Add a caption-like paragraph before each leftover image
            doc.add_paragraph("")  # spacer
            _add_image(doc, img)

    # Save document
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output))

    print(f"[OK] Created DOCX: {output_path}")
    print(f"  - Title: {title}")
    print(f"  - Paragraphs: {len(doc.paragraphs)}")
    print(f"  - Images: {len(images)}")

    return str(output)


def create_xlsx(
    data_path: Optional[str] = None,
    sheet_name: str = "Sheet1",
    output_path: str = "output.xlsx",
    title: Optional[str] = None,
) -> str:
    """
    Create an Excel spreadsheet with data from CSV or arrays.

    Args:
        data_path: Path to CSV file
        sheet_name: Name of the sheet
        output_path: Output file path
        title: Optional title for the spreadsheet

    Returns:
        Path to created spreadsheet
    """
    # Ensure output path has .xlsx extension
    if not output_path.endswith('.xlsx'):
        output_path = f"{output_path}.xlsx"

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name

    # Add title if provided
    if title:
        ws['A1'] = title
        ws['A1'].font = Font(size=16, bold=True, name=_get_system_cjk_font())
        start_row = 2
    else:
        start_row = 1

    # Load data from CSV
    if data_path:
        csv_path = Path(data_path)
        if csv_path.exists():
            import csv
            with open(csv_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                data = list(reader)

            # Write data to sheet
            for row_idx, row in enumerate(data, start=start_row):
                for col_idx, value in enumerate(row, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)

                    # Format header row
                    if row_idx == start_row:
                        cell.font = Font(bold=True, name=_get_system_cjk_font())
                        cell.fill = PatternFill(start_color="DDDDDD", end_color="DDDDDD", fill_type="solid")

            # Auto-adjust column widths
            for col_idx in range(1, len(data[0]) + 1):
                column_letter = get_column_letter(col_idx)
                ws.column_dimensions[column_letter].width = 15

            print(f"[OK] Loaded {len(data)} rows from {data_path}")
        else:
            print(f"Warning: Data file not found: {data_path}")
            # Add sample data
            ws.append(["Column1", "Column2", "Column3"])
            ws.append(["Data1", "Data2", "Data3"])
    else:
        # Add sample data
        ws.append(["Column1", "Column2", "Column3"])
        ws.append(["Data1", "Data2", "Data3"])

    # Save spreadsheet
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output))

    print(f"[OK] Created XLSX: {output_path}")
    print(f"  - Sheet: {sheet_name}")
    print(f"  - Rows: {ws.max_row}")

    return str(output)


def create_pptx(
    title: str,
    content: Optional[str] = None,
    image_path: Optional[str] = None,
    slides: int = 1,
    output_path: str = "output.pptx",
) -> str:
    """
    Create a PowerPoint presentation with title, content, and images.

    Args:
        title: Presentation title
        content: Content text (can be comma-separated for multiple slides)
        image_path: Optional path to image file
        slides: Number of slides
        output_path: Output file path

    Returns:
        Path to created presentation
    """
    # Ensure output path has .pptx extension
    if not output_path.endswith('.pptx'):
        output_path = f"{output_path}.pptx"

    prs = Presentation()

    # Slide 1: Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout
    title_shape = title_slide.shapes.title
    subtitle_shape = title_slide.placeholders[1]

    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = _get_system_cjk_font()
    subtitle_shape.text = "Generated by Doc Creator"

    # Parse content for multiple slides
    content_items = []
    if content:
        if ',' in content:
            content_items = [c.strip() for c in content.split(',')]
        else:
            content_items = [content]

    # Create content slides
    for i in range(min(slides - 1, len(content_items) if content_items else slides - 1)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout

        if i < len(content_items):
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.text = content_items[i]
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = _get_system_cjk_font()

            # Add image if provided and on first content slide
            if image_path and i == 0:
                img_path = Path(image_path)
                if img_path.exists():
                    try:
                        # Add image to slide
                        slide.shapes.add_picture(
                            str(img_path),
                            Inches(0.5),
                            Inches(2),
                            width=Inches(9)
                        )
                    except Exception as e:
                        print(f"Warning: Could not insert image: {e}")
        else:
            # Placeholder slide
            title_shape = slide.shapes.title
            title_shape.text = f"Slide {i + 2}"

    # Save presentation
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    print(f"[OK] Created PPTX: {output_path}")
    print(f"  - Slides: {len(prs.slides)}")
    print(f"  - Title: {title}")

    return str(output)


def main():
    parser = argparse.ArgumentParser(
        description="Create DOCX/XLSX/PPTX documents",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Word document with image
  python doc_creator.py --type docx --title "Report" --content "Summary..." --image chart.png

  # Excel from CSV
  python doc_creator.py --type xlsx --data data.csv --sheet-name "Sales"

  # PowerPoint presentation
  python doc_creator.py --type pptx --title "Review" --content "Agenda,Metrics,Next Steps" --slides 4
        """
    )

    parser.add_argument(
        "--type",
        required=True,
        choices=["docx", "xlsx", "pptx"],
        help="Document type"
    )

    parser.add_argument(
        "--title",
        help="Document title"
    )

    parser.add_argument(
        "--content",
        help="Document content (for DOCX/PPTX)"
    )

    parser.add_argument(
        "--image",
        action="append",
        help="Path to image file (can specify multiple times, for DOCX/PPTX)"
    )

    parser.add_argument(
        "--font",
        help="CJK font name (default: auto-detect by OS)"
    )

    parser.add_argument(
        "--data",
        help="Path to CSV data file (for XLSX)"
    )

    parser.add_argument(
        "--csv",
        dest="data",
        help="Alias for --data"
    )

    parser.add_argument(
        "--sheet-name",
        default="Sheet1",
        help="Sheet name for XLSX (default: Sheet1)"
    )

    parser.add_argument(
        "--slides",
        type=int,
        default=1,
        help="Number of slides for PPTX (default: 1)"
    )

    parser.add_argument(
        "--output",
        default="output",
        help="Output file path without extension (default: output)"
    )

    args = parser.parse_args()

    # Build output path with extension (if not already present)
    if args.output.endswith(f".{args.type}"):
        output_path = args.output
    else:
        output_path = f"{args.output}.{args.type}"

    try:
        if args.type == "docx":
            create_docx(
                title=args.title or "Document",
                content=args.content or "",
                image_paths=args.image or None,
                output_path=output_path,
                font_name=args.font,
            )
        elif args.type == "xlsx":
            create_xlsx(
                data_path=args.data,
                sheet_name=args.sheet_name,
                output_path=output_path,
                title=args.title,
            )
        elif args.type == "pptx":
            create_pptx(
                title=args.title or "Presentation",
                content=args.content,
                image_path=args.image,
                slides=args.slides,
                output_path=output_path,
            )
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
